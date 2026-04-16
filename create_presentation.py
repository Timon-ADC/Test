from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation('Template.pptx')

# Remove all existing template slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# ---- Helpers ----
def clear_and_set(placeholder, text, font_size=None, bold=None, color="000000"):
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_para(tf, text, size=11, bold=False, color="000000", space_before=None):
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    if space_before:
        p.space_before = Pt(space_before)
    return p


def set_first_para(tf, text, size=16, bold=True, color="000000"):
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


# ============================================================
# SLIDE 1: Title
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[0])
clear_and_set(s1.placeholders[0], "Collective Week\nRapid Prototyping", font_size=36, bold=True)
clear_and_set(s1.placeholders[1], "Timon Domela Nieuwenhuis Nyegaard — ADC Consulting", font_size=14)

# ============================================================
# SLIDE 2: Expectations vs. Reality
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[16])
clear_and_set(s2.placeholders[0], "Expectations vs. Reality", font_size=28, bold=True)

tf_l = s2.placeholders[1].text_frame
tf_l.clear()
set_first_para(tf_l, "Going In", size=18, bold=True)
add_para(tf_l, "", size=6)
add_para(tf_l, "Focus mainly on rapid prototyping techniques", size=12, space_before=8)
add_para(tf_l, "", size=4)
add_para(tf_l, "Classic hackathon: long hours, lots of code, heads down", size=12, space_before=8)
add_para(tf_l, "", size=4)
add_para(tf_l, "Meet colleagues, mostly within my usual circles", size=12, space_before=8)

tf_r = s2.placeholders[2].text_frame
tf_r.clear()
set_first_para(tf_r, "What Actually Happened", size=18, bold=True)
add_para(tf_r, "", size=6)
add_para(tf_r, "A true rapid prototyping lab with fast build-test-iterate cycles", size=12, space_before=8)
add_para(tf_r, "", size=4)
add_para(tf_r, "Hackathon mixed hard work with playful energy — including karaoke", size=12, space_before=8)
add_para(tf_r, "", size=4)
add_para(tf_r, "Connected across roles and personalities, not just my usual team", size=12, space_before=8)

# ============================================================
# SLIDE 3: Highlights
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[5])
clear_and_set(s3.placeholders[0], "Hackathon Energy\n& Karaoke Vibes", font_size=28, bold=True)
clear_and_set(s3.placeholders[10], "The week's most memorable moments", font_size=14, bold=True)

tf_m = s3.placeholders[11].text_frame
tf_m.clear()
set_first_para(tf_m, "The Hackathon", size=16, bold=True)
add_para(tf_m, "Pushing from idea to working prototype under time pressure showed how much you can achieve when you stop over-planning and just build.", size=12, space_before=6)
add_para(tf_m, "", size=10)
add_para(tf_m, "Karaoke Night", size=16, bold=True, space_before=12)
add_para(tf_m, "Completely different energy, but it broke the ice, flattened hierarchies, and made collaboration the next day noticeably smoother.", size=12, space_before=6)
add_para(tf_m, "", size=10)
add_para(tf_m, "The Atmosphere", size=16, bold=True, space_before=12)
add_para(tf_m, "A mix of focus and fun that made experimentation feel safe — it was okay if something broke, as long as we learned.", size=12, space_before=6)

# ============================================================
# SLIDE 4: What I Learned (3 columns)
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[18])
clear_and_set(s4.placeholders[0], '"Show, Don\'t Tell" — My Key Learnings', font_size=28, bold=True)

tf1 = s4.placeholders[14].text_frame
tf1.clear()
set_first_para(tf1, "On the Craft", size=16, bold=True)
add_para(tf1, "", size=6)
add_para(tf1, "Instead of explaining ideas endlessly, build a quick version and let people react to something tangible.", size=11, space_before=6)
add_para(tf1, "", size=6)
add_para(tf1, '"Show, don\'t tell" made feedback concrete and immediate.', size=11, space_before=6)

tf2 = s4.placeholders[15].text_frame
tf2.clear()
set_first_para(tf2, "On Collaboration", size=16, bold=True)
add_para(tf2, "", size=6)
add_para(tf2, "Prototyping together created alignment much faster than long discussions or slide decks.", size=11, space_before=6)
add_para(tf2, "", size=6)
add_para(tf2, "People rallied around what they could see and test, not abstract descriptions.", size=11, space_before=6)

tf3 = s4.placeholders[16].text_frame
tf3.clear()
set_first_para(tf3, "On Myself", size=16, bold=True)
add_para(tf3, "", size=6)
add_para(tf3, 'I realised I can "vibe code" too — jumping into code intuitively, exploring and adjusting on the fly.', size=11, space_before=6)
add_para(tf3, "", size=6)
add_para(tf3, "I don't need everything perfectly specified before I start building.", size=11, space_before=6)

# ============================================================
# SLIDE 5: Closing
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[25])
clear_and_set(s5.placeholders[0], "How This Changes\nMy Work", font_size=28, bold=True)

tf5 = s5.placeholders[1000].text_frame
tf5.clear()
set_first_para(tf5, "What I'll Do Differently", size=16, bold=True)
add_para(tf5, "", size=6)
add_para(tf5, "Prototype earlier and more roughly to get feedback faster", size=11, space_before=6)
add_para(tf5, "", size=4)
add_para(tf5, 'Use "show don\'t tell" in client and internal work — demos over descriptions', size=11, space_before=6)
add_para(tf5, "", size=4)
add_para(tf5, "Lean into vibe coding as a way to explore ideas quickly", size=11, space_before=6)
add_para(tf5, "", size=10)
add_para(tf5, "For Next Time", size=16, bold=True, space_before=12)
add_para(tf5, "", size=6)
add_para(tf5, "Location in Spain — a change of environment to boost energy and creativity even more", size=11, space_before=6)
add_para(tf5, "", size=4)
add_para(tf5, "Keep combining structured sprints with playful moments to maintain openness", size=11, space_before=6)
add_para(tf5, "", size=10)
add_para(tf5, "Speed, play, and prototypes can be just as important as plans, specs, and slide decks.", size=12, bold=True, space_before=12)

# ============================================================
# SAVE
# ============================================================
prs.save('collective_week_presentation.pptx')
print("Done! File saved as collective_week_presentation.pptx")
